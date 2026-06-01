# Copyright (c) 2024-2026 T1 Agentics LLC
# SPDX-License-Identifier: Apache-2.0

"""
Knowledge Base Document Processor

AI-powered document processing for the Company Best Practices Database.
Handles:
- Document upload and text extraction (PDF, DOCX, TXT)
- AI analysis to extract rules, procedures, and key points
- Automatic categorization and tagging
- SOAR playbook import and conversion
"""

import json
import logging
import uuid
import os
import base64
from datetime import datetime
from typing import Dict, Any, Optional, List

logger = logging.getLogger(__name__)


class KBDocumentProcessor:
    """
    AI-powered document processor for the knowledge base.

    Features:
    - Extract text from uploaded documents
    - Use AI to analyze and structure content
    - Extract actionable rules and procedures
    - Auto-categorize and tag content
    - Support SOAR playbook conversion
    """

    def __init__(self):
        self.enabled = True
        self.model = os.getenv('AI_KB_MODEL', 'claude-sonnet-4-20250514')
        # Supported file types based on format spec:
        # required: md, txt, json
        # recommended: pdf, csv, yaml
        # optional: html, docx, pptx
        self.supported_types = ['md', 'txt', 'json', 'pdf', 'csv', 'yaml', 'yml', 'html', 'docx', 'pptx']

    async def process_document(
        self,
        filename: str,
        content: bytes,
        file_type: str,
        uploaded_by: str = 'system'
    ) -> Dict[str, Any]:
        """
        Process an uploaded document and extract knowledge base entries.

        Args:
            filename: Original filename
            content: File content as bytes
            file_type: File extension/type
            uploaded_by: Username of uploader

        Returns:
            Processing result with extracted entries
        """
        try:
            from services.postgres_db import postgres_db

            upload_id = f"upload-{uuid.uuid4().hex[:8]}"

            # Create upload record
            if postgres_db.connected:
                async with postgres_db.tenant_acquire() as conn:
                    await conn.execute('''
                        INSERT INTO kb_document_uploads (
                            upload_id, filename, file_type, file_size,
                            status, uploaded_by
                        ) VALUES ($1, $2, $3, $4, 'processing', $5)
                    ''', upload_id, filename, file_type, len(content), uploaded_by)

            logger.info(f"Processing document: {filename} ({file_type}, {len(content)} bytes)")

            # Extract text from document
            extracted_text = await self._extract_text(content, file_type)

            if not extracted_text:
                await self._update_upload_status(upload_id, 'failed', error="Could not extract text from document")
                return {"error": "Could not extract text from document", "upload_id": upload_id}

            # Analyze with AI
            ai_analysis = await self._analyze_document(extracted_text, filename, file_type)

            if ai_analysis.get('error'):
                await self._update_upload_status(upload_id, 'failed', error=ai_analysis['error'])
                return {"error": ai_analysis['error'], "upload_id": upload_id}

            # Create knowledge base entries from analysis
            created_entries = await self._create_entries_from_analysis(
                ai_analysis,
                filename,
                file_type,
                uploaded_by
            )

            # Update upload record with results
            if postgres_db.connected:
                async with postgres_db.tenant_acquire() as conn:
                    # kb_document_uploads schema only has status / completed_at /
                    # resulting_kb_ids on the write side. Extracted text +
                    # AI analysis live in the created knowledge_base rows.
                    await conn.execute('''
                        UPDATE kb_document_uploads
                        SET status = 'completed',
                            completed_at = CURRENT_TIMESTAMP,
                            resulting_kb_ids = $1
                        WHERE upload_id = $2
                    ''',
                        [e['kb_id'] for e in created_entries],
                        upload_id
                    )

            return {
                "upload_id": upload_id,
                "status": "completed",
                "filename": filename,
                "extracted_text_length": len(extracted_text),
                "ai_analysis": ai_analysis,
                "created_entries": created_entries,
                "entries_count": len(created_entries)
            }

        except Exception as e:
            logger.error(f"Document processing failed: {e}")
            return {"error": str(e)}

    async def _extract_text(self, content: bytes, file_type: str) -> Optional[str]:
        """Extract text from document based on type."""

        file_type = file_type.lower().lstrip('.')

        if file_type in ['txt', 'md']:
            # Plain text
            try:
                return content.decode('utf-8')
            except UnicodeDecodeError:
                return content.decode('latin-1')

        elif file_type in ['json']:
            # JSON - pretty print
            try:
                data = json.loads(content.decode('utf-8'))
                return json.dumps(data, indent=2)
            except:
                return content.decode('utf-8')

        elif file_type in ['yaml', 'yml']:
            # YAML - return as-is
            try:
                return content.decode('utf-8')
            except:
                return content.decode('latin-1')

        elif file_type == 'pdf':
            # PDF extraction (requires pypdf or similar)
            try:
                import io
                try:
                    from pypdf import PdfReader
                except ImportError:
                    from PyPDF2 import PdfReader

                reader = PdfReader(io.BytesIO(content))
                text_parts = []
                for page in reader.pages:
                    text = page.extract_text()
                    if text:
                        text_parts.append(text)
                return '\n\n'.join(text_parts)
            except ImportError:
                logger.warning("PDF library not installed. Install pypdf or PyPDF2")
                return None
            except Exception as e:
                logger.error(f"PDF extraction failed: {e}")
                return None

        elif file_type == 'docx':
            # DOCX extraction (requires python-docx)
            try:
                import io
                from docx import Document

                doc = Document(io.BytesIO(content))
                text_parts = []
                for paragraph in doc.paragraphs:
                    if paragraph.text.strip():
                        text_parts.append(paragraph.text)
                return '\n\n'.join(text_parts)
            except ImportError:
                logger.warning("python-docx not installed. Install python-docx")
                return None
            except Exception as e:
                logger.error(f"DOCX extraction failed: {e}")
                return None

        elif file_type == 'csv':
            # CSV - convert to markdown table
            try:
                import csv
                import io

                text = content.decode('utf-8')
                reader = csv.reader(io.StringIO(text))
                rows = list(reader)

                if not rows:
                    return text

                # Convert to markdown table
                md_lines = []
                header = rows[0]
                md_lines.append('| ' + ' | '.join(header) + ' |')
                md_lines.append('| ' + ' | '.join(['---'] * len(header)) + ' |')

                for row in rows[1:]:
                    # Pad row if needed
                    while len(row) < len(header):
                        row.append('')
                    md_lines.append('| ' + ' | '.join(row[:len(header)]) + ' |')

                return '\n'.join(md_lines)
            except Exception as e:
                logger.error(f"CSV extraction failed: {e}")
                # Fall back to raw text
                try:
                    return content.decode('utf-8')
                except:
                    return None

        elif file_type == 'html':
            # HTML - strip tags and extract text
            try:
                import re

                text = content.decode('utf-8')

                # Remove script and style elements
                text = re.sub(r'<script[^>]*>.*?</script>', '', text, flags=re.DOTALL | re.IGNORECASE)
                text = re.sub(r'<style[^>]*>.*?</style>', '', text, flags=re.DOTALL | re.IGNORECASE)

                # Convert common HTML to markdown-ish
                text = re.sub(r'<h1[^>]*>(.*?)</h1>', r'\n# \1\n', text, flags=re.IGNORECASE)
                text = re.sub(r'<h2[^>]*>(.*?)</h2>', r'\n## \1\n', text, flags=re.IGNORECASE)
                text = re.sub(r'<h3[^>]*>(.*?)</h3>', r'\n### \1\n', text, flags=re.IGNORECASE)
                text = re.sub(r'<p[^>]*>(.*?)</p>', r'\n\1\n', text, flags=re.IGNORECASE | re.DOTALL)
                text = re.sub(r'<li[^>]*>(.*?)</li>', r'- \1', text, flags=re.IGNORECASE)
                text = re.sub(r'<br\s*/?>', '\n', text, flags=re.IGNORECASE)

                # Remove remaining tags
                text = re.sub(r'<[^>]+>', '', text)

                # Decode HTML entities
                text = text.replace('&nbsp;', ' ')
                text = text.replace('&amp;', '&')
                text = text.replace('&lt;', '<')
                text = text.replace('&gt;', '>')
                text = text.replace('&quot;', '"')

                # Clean up whitespace
                text = re.sub(r'\n\s*\n', '\n\n', text)
                text = text.strip()

                return text
            except Exception as e:
                logger.error(f"HTML extraction failed: {e}")
                return None

        elif file_type == 'pptx':
            # PPTX extraction (requires python-pptx)
            try:
                import io
                from pptx import Presentation

                prs = Presentation(io.BytesIO(content))
                text_parts = []

                for slide_num, slide in enumerate(prs.slides, 1):
                    slide_text = [f"## Slide {slide_num}"]

                    for shape in slide.shapes:
                        if hasattr(shape, "text") and shape.text.strip():
                            slide_text.append(shape.text.strip())

                    if len(slide_text) > 1:  # More than just the slide header
                        text_parts.append('\n'.join(slide_text))

                return '\n\n'.join(text_parts)
            except ImportError:
                logger.warning("python-pptx not installed. Install python-pptx")
                return None
            except Exception as e:
                logger.error(f"PPTX extraction failed: {e}")
                return None

        else:
            # Unknown type - try as text
            try:
                return content.decode('utf-8')
            except:
                return None

    def _parse_sop_frontmatter(self, text: str) -> tuple[Dict[str, Any], str]:
        """
        Parse YAML frontmatter from an SOP document.

        SOPs should have structured metadata in YAML format between --- delimiters:
        ---
        sop_id: SOP-XXX-001
        title: Title
        scope: tier1|tier2|tier3
        applies_to: [alert_type, entity_type]
        priority: low|medium|high
        confidence_threshold: 0.7
        allowed_actions: [close, escalate, block, isolate, notify]
        required_conditions: [condition1, condition2]
        ---

        Returns:
            Tuple of (metadata dict, body text without frontmatter)
        """
        import yaml

        metadata = {}
        body = text

        # Check for YAML frontmatter (--- at start, --- to close)
        if text.strip().startswith('---'):
            parts = text.split('---', 2)
            if len(parts) >= 3:
                try:
                    # parts[0] is empty (before first ---), parts[1] is YAML, parts[2] is body
                    yaml_content = parts[1].strip()
                    metadata = yaml.safe_load(yaml_content) or {}
                    body = parts[2].strip()

                    logger.info(f"Parsed SOP frontmatter: sop_id={metadata.get('sop_id')}")
                except yaml.YAMLError as e:
                    logger.warning(f"Failed to parse YAML frontmatter: {e}")
                except Exception as e:
                    logger.warning(f"Error parsing frontmatter: {e}")

        return metadata, body

    def _convert_frontmatter_to_kb_fields(self, metadata: Dict[str, Any]) -> Dict[str, Any]:
        """
        Convert SOP frontmatter metadata to knowledge base entry fields.

        Maps:
        - sop_id -> kb_id prefix
        - scope -> tier filter
        - applies_to -> incident_types, ioc_types
        - priority -> priority number
        - confidence_threshold -> stored in ai_extracted_rules
        - allowed_actions -> stored in ai_extracted_rules
        - required_conditions -> stored in ai_extracted_rules
        """
        result = {}

        # Map priority text to number (lower = higher priority)
        priority_map = {'critical': 10, 'high': 25, 'medium': 50, 'low': 75}
        if metadata.get('priority'):
            result['priority'] = priority_map.get(metadata['priority'].lower(), 50)

        # Map scope to severity filter
        scope = metadata.get('scope', '').lower()
        if scope == 'tier1':
            result['severity_filter'] = ['low', 'medium', 'high', 'critical']
        elif scope == 'tier2':
            result['severity_filter'] = ['medium', 'high', 'critical']
        elif scope == 'tier3':
            result['severity_filter'] = ['high', 'critical']

        # Map applies_to to incident_types and ioc_types
        applies_to = metadata.get('applies_to', [])
        if isinstance(applies_to, str):
            applies_to = [applies_to]

        ioc_keywords = ['ip', 'domain', 'hash', 'url', 'email', 'file']
        incident_keywords = ['phishing', 'malware', 'ransomware', 'authentication',
                           'network', 'email', 'data', 'insider', 'generic']

        ioc_types = []
        incident_types = []

        for item in applies_to:
            item_lower = item.lower()
            if item_lower in ioc_keywords:
                ioc_types.append(item_lower)
            if any(kw in item_lower for kw in incident_keywords):
                incident_types.append(item_lower)
            # Also add as tag
            if 'tags' not in result:
                result['tags'] = []
            result['tags'].append(item_lower)

        if ioc_types:
            result['ioc_types'] = ioc_types
        if incident_types:
            result['incident_types'] = incident_types

        # Store SOP-specific metadata as extracted rules
        sop_rules = []

        if metadata.get('sop_id'):
            sop_rules.append(f"SOP ID: {metadata['sop_id']}")

        if metadata.get('confidence_threshold'):
            sop_rules.append(f"Confidence threshold: {metadata['confidence_threshold']}")

        if metadata.get('allowed_actions'):
            actions = metadata['allowed_actions']
            if isinstance(actions, list):
                sop_rules.append(f"Allowed actions: {', '.join(actions)}")

        if metadata.get('required_conditions'):
            conditions = metadata['required_conditions']
            if isinstance(conditions, list):
                sop_rules.append(f"Required conditions: {', '.join(conditions)}")

        if sop_rules:
            result['sop_metadata_rules'] = sop_rules

        return result

    def _merge_frontmatter_into_entry(
        self,
        entry: Dict[str, Any],
        frontmatter: Dict[str, Any],
        frontmatter_fields: Dict[str, Any]
    ) -> None:
        """
        Merge SOP frontmatter metadata into a KB entry.

        Frontmatter fields take precedence over AI-extracted fields for structured data.
        This ensures SOPs maintain their explicit metadata while enriching with AI analysis.
        """
        # Update title with SOP ID prefix
        if frontmatter.get('sop_id'):
            original_title = entry.get('title', 'Untitled')
            fm_title = frontmatter.get('title', original_title)
            entry['title'] = f"{frontmatter['sop_id']} — {fm_title}"

        # Merge priority (frontmatter takes precedence)
        if 'priority' in frontmatter_fields:
            entry['priority'] = frontmatter_fields['priority']

        # Merge severity filter
        if 'severity_filter' in frontmatter_fields:
            entry['severity_filter'] = frontmatter_fields['severity_filter']

        # Merge IOC types (union of frontmatter and AI)
        if 'ioc_types' in frontmatter_fields:
            existing = entry.get('ioc_types', [])
            entry['ioc_types'] = list(set(existing + frontmatter_fields['ioc_types']))

        # Merge incident types (union of frontmatter and AI)
        if 'incident_types' in frontmatter_fields:
            existing = entry.get('incident_types', [])
            entry['incident_types'] = list(set(existing + frontmatter_fields['incident_types']))

        # Merge tags (union of frontmatter and AI)
        if 'tags' in frontmatter_fields:
            existing = entry.get('tags', [])
            entry['tags'] = list(set(existing + frontmatter_fields['tags']))

        # Add SOP metadata rules to extracted rules
        if 'sop_metadata_rules' in frontmatter_fields:
            existing_rules = entry.get('extracted_rules', [])
            # Prepend SOP metadata rules
            entry['extracted_rules'] = frontmatter_fields['sop_metadata_rules'] + existing_rules

        # Store the raw frontmatter for agent access
        entry['sop_frontmatter'] = frontmatter

    async def _analyze_document(
        self,
        text: str,
        filename: str,
        file_type: str
    ) -> Dict[str, Any]:
        """Use AI to analyze document and extract structured content."""

        # First, check for SOP frontmatter (structured metadata)
        frontmatter, body_text = self._parse_sop_frontmatter(text)
        frontmatter_fields = {}

        if frontmatter.get('sop_id'):
            # This is a properly formatted SOP - extract metadata
            frontmatter_fields = self._convert_frontmatter_to_kb_fields(frontmatter)
            logger.info(f"Processing SOP with frontmatter: {frontmatter.get('sop_id')}")

        api_key = os.getenv('ANTHROPIC_API_KEY')
        if not api_key:
            logger.warning("ANTHROPIC_API_KEY not set, using mock analysis")
            result = self._mock_analysis(text, filename)
            # Merge frontmatter data
            if frontmatter_fields and result.get('entries'):
                for entry in result['entries']:
                    self._merge_frontmatter_into_entry(entry, frontmatter, frontmatter_fields)
            return result

        # Build the prompt (use original text to include frontmatter context)
        prompt = self._build_analysis_prompt(text, filename, file_type)

        try:
            import aiohttp

            headers = {
                "x-api-key": api_key,
                "anthropic-version": "2023-06-01",
                "content-type": "application/json"
            }

            payload = {
                "model": self.model,
                "max_tokens": 4000,
                "messages": [
                    {"role": "user", "content": prompt}
                ]
            }

            async with aiohttp.ClientSession() as session:
                async with session.post(
                    "https://api.anthropic.com/v1/messages",
                    headers=headers,
                    json=payload,
                    timeout=aiohttp.ClientTimeout(total=60)
                ) as response:
                    if response.status != 200:
                        error_text = await response.text()
                        logger.error(f"AI API error: {response.status} - {error_text}")
                        result = self._mock_analysis(text, filename)
                        # Still merge frontmatter on fallback
                        if frontmatter_fields and result.get('entries'):
                            for entry in result['entries']:
                                self._merge_frontmatter_into_entry(entry, frontmatter, frontmatter_fields)
                        return result

                    data = await response.json()
                    response_text = data['content'][0]['text']

                    result = self._parse_analysis_response(response_text)

                    # Merge frontmatter data into AI analysis results
                    if frontmatter_fields and result.get('entries'):
                        for entry in result['entries']:
                            self._merge_frontmatter_into_entry(entry, frontmatter, frontmatter_fields)

                    return result

        except Exception as e:
            logger.error(f"AI analysis failed: {e}")
            return self._mock_analysis(text, filename)

    def _build_analysis_prompt(self, text: str, filename: str, file_type: str) -> str:
        """Build the document analysis prompt."""

        # Truncate very long documents
        max_chars = 30000
        if len(text) > max_chars:
            text = text[:max_chars] + "\n\n[Document truncated...]"

        return f"""You are a SOC documentation analyst. Analyze this document and extract structured knowledge base entries for a Security Operations Center.

DOCUMENT: {filename} ({file_type})

CONTENT:
{text}

Analyze this document and extract one or more knowledge base entries. For each entry, identify:
1. A clear, descriptive title
2. The type of content (sop, playbook, escalation, compliance, policy, procedure, runbook, handling_rule)
3. The category (incident_response, threat_detection, malware_analysis, phishing, data_loss, insider_threat, network_security, endpoint_security, cloud_security, identity_access, compliance, escalation, communication, documentation, general)
4. Relevant tags
5. Which severities it applies to (low, medium, high, critical)
6. Which incident types it applies to
7. Which IOC types it's relevant for (ip, domain, hash, url, email)
8. MITRE ATT&CK techniques if applicable
9. Compliance frameworks if applicable (NIST, SOC2, ISO27001, PCI-DSS, HIPAA)
10. A priority (1-1000, lower = higher priority)
11. An AI summary (2-3 sentences)
12. Extracted actionable rules (list of specific steps or rules)

Respond with JSON only (no markdown, no preamble):
{{
  "document_summary": "Brief summary of the entire document",
  "entries": [
    {{
      "title": "Entry title",
      "content_type": "sop|playbook|escalation|compliance|policy|procedure|runbook|handling_rule",
      "category": "category from list above",
      "content": "Full content for this entry - include all relevant procedures, steps, and details",
      "tags": ["tag1", "tag2"],
      "severity_filter": ["low", "medium", "high", "critical"],
      "incident_types": ["phishing", "malware", "data_breach", etc],
      "ioc_types": ["ip", "domain", "hash"],
      "mitre_techniques": ["T1566", "T1059"],
      "compliance_frameworks": ["NIST CSF", "SOC2"],
      "priority": 100,
      "ai_summary": "2-3 sentence summary",
      "extracted_rules": [
        "Rule or step 1",
        "Rule or step 2"
      ]
    }}
  ]
}}

If the document contains multiple distinct procedures or policies, create separate entries for each.
If it's a single coherent document, create one comprehensive entry."""

    def _parse_analysis_response(self, response_text: str) -> Dict[str, Any]:
        """Parse the AI analysis response."""
        try:
            # Remove markdown code blocks if present
            text = response_text.strip()
            if text.startswith('```'):
                lines = text.split('\n')
                text = '\n'.join(lines[1:-1])

            result = json.loads(text)
            result['ai_model'] = self.model
            result['timestamp'] = datetime.utcnow().isoformat()
            return result

        except json.JSONDecodeError as e:
            logger.error(f"Failed to parse AI analysis response: {e}")
            return {
                "error": "Failed to parse AI response",
                "raw_response": response_text[:1000]
            }

    def _mock_analysis(self, text: str, filename: str) -> Dict[str, Any]:
        """Generate mock analysis when AI is unavailable."""

        # Simple heuristic extraction
        lines = text.split('\n')
        title = lines[0][:100] if lines else filename

        return {
            "document_summary": f"Document: {filename}",
            "entries": [
                {
                    "title": title,
                    "content_type": "sop",
                    "category": "general",
                    "content": text[:5000],
                    "tags": ["imported"],
                    "severity_filter": [],
                    "incident_types": [],
                    "ioc_types": [],
                    "mitre_techniques": [],
                    "compliance_frameworks": [],
                    "priority": 100,
                    "ai_summary": f"Imported from {filename}. Manual review recommended.",
                    "extracted_rules": []
                }
            ],
            "ai_model": "mock",
            "timestamp": datetime.utcnow().isoformat()
        }

    async def _create_entries_from_analysis(
        self,
        analysis: Dict[str, Any],
        source_filename: str,
        source_type: str,
        created_by: str
    ) -> List[Dict[str, Any]]:
        """Create knowledge base entries from AI analysis."""

        from services.knowledge_base_service import get_knowledge_base_service

        kb_service = get_knowledge_base_service()
        created_entries = []

        entries = analysis.get('entries', [])

        for entry_data in entries:
            try:
                result = await kb_service.create_entry(
                    title=entry_data.get('title', 'Untitled'),
                    content=entry_data.get('content', ''),
                    content_type=entry_data.get('content_type', 'sop'),
                    category=entry_data.get('category'),
                    tags=entry_data.get('tags', []),
                    severity_filter=entry_data.get('severity_filter', []),
                    incident_types=entry_data.get('incident_types', []),
                    ioc_types=entry_data.get('ioc_types', []),
                    mitre_techniques=entry_data.get('mitre_techniques', []),
                    compliance_frameworks=entry_data.get('compliance_frameworks', []),
                    priority=entry_data.get('priority', 100),
                    created_by=created_by,
                    source_document_name=source_filename,
                    source_document_type=source_type
                )

                if not result.get('error'):
                    # Build the AI extracted rules with SOP frontmatter if present
                    extracted_rules_data = {
                        'rules': entry_data.get('extracted_rules', [])
                    }

                    # Include SOP frontmatter for agent access
                    if entry_data.get('sop_frontmatter'):
                        fm = entry_data['sop_frontmatter']
                        extracted_rules_data['sop_metadata'] = {
                            'sop_id': fm.get('sop_id'),
                            'scope': fm.get('scope'),
                            'applies_to': fm.get('applies_to', []),
                            'priority': fm.get('priority'),
                            'confidence_threshold': fm.get('confidence_threshold'),
                            'allowed_actions': fm.get('allowed_actions', []),
                            'required_conditions': fm.get('required_conditions', [])
                        }

                    # Update with AI analysis
                    await kb_service.update_entry(
                        result['kb_id'],
                        {
                            'ai_processed': True,
                            'ai_summary': entry_data.get('ai_summary'),
                            'ai_extracted_rules': extracted_rules_data
                        },
                        updated_by='ai_processor'
                    )
                    created_entries.append(result)
                else:
                    logger.error(f"Failed to create KB entry: {result['error']}")

            except Exception as e:
                logger.error(f"Error creating entry from analysis: {e}")

        return created_entries

    async def _update_upload_status(
        self,
        upload_id: str,
        status: str,
        error: Optional[str] = None
    ):
        """Update upload record status."""
        try:
            from services.postgres_db import postgres_db

            if postgres_db.connected:
                async with postgres_db.tenant_acquire() as conn:
                    await conn.execute('''
                        UPDATE kb_document_uploads
                        SET status = $1, error_message = $2, completed_at = CURRENT_TIMESTAMP
                        WHERE upload_id = $3
                    ''', status, error, upload_id)
        except Exception as e:
            logger.error(f"Failed to update upload status: {e}")

    async def get_upload_status(self, upload_id: str) -> Optional[Dict[str, Any]]:
        """Get status of a document upload/processing job."""
        try:
            from services.postgres_db import postgres_db

            if not postgres_db.connected:
                return None

            async with postgres_db.tenant_acquire() as conn:
                row = await conn.fetchrow(
                    'SELECT * FROM kb_document_uploads WHERE upload_id = $1',
                    upload_id
                )

                if not row:
                    return None

                result = dict(row)

                # Convert UUID to string
                if result.get('id'):
                    result['id'] = str(result['id'])

                # Convert datetimes
                for field in ['created_at', 'completed_at']:
                    if result.get(field):
                        result[field] = result[field].isoformat()

                # Parse AI analysis if present
                if result.get('ai_analysis'):
                    if isinstance(result['ai_analysis'], str):
                        try:
                            result['ai_analysis'] = json.loads(result['ai_analysis'])
                        except:
                            pass

                return result

        except Exception as e:
            logger.error(f"Failed to get upload status: {e}")
            return None

    async def list_uploads(
        self,
        status: Optional[str] = None,
        limit: int = 50
    ) -> List[Dict[str, Any]]:
        """List document uploads."""
        try:
            from services.postgres_db import postgres_db

            if not postgres_db.connected:
                return []

            async with postgres_db.tenant_acquire() as conn:
                if status:
                    rows = await conn.fetch('''
                        SELECT upload_id, filename, file_type, file_size,
                               status, error_message, resulting_kb_ids,
                               uploaded_by, created_at, completed_at
                        FROM kb_document_uploads
                        WHERE status = $1
                        ORDER BY created_at DESC
                        LIMIT $2
                    ''', status, limit)
                else:
                    rows = await conn.fetch('''
                        SELECT upload_id, filename, file_type, file_size,
                               status, error_message, resulting_kb_ids,
                               uploaded_by, created_at, completed_at
                        FROM kb_document_uploads
                        ORDER BY created_at DESC
                        LIMIT $1
                    ''', limit)

                results = []
                for row in rows:
                    result = dict(row)
                    for field in ['created_at', 'completed_at']:
                        if result.get(field):
                            result[field] = result[field].isoformat()
                    results.append(result)

                return results

        except Exception as e:
            logger.error(f"Failed to list uploads: {e}")
            return []

    async def process_soar_playbook(
        self,
        playbook_content: str,
        playbook_format: str,  # palo_xsoar, swimlane, etc.
        uploaded_by: str = 'system'
    ) -> Dict[str, Any]:
        """
        Process and convert a SOAR playbook to knowledge base entries.

        This is a specialized processor for importing playbooks from
        various SOAR platforms and converting them into intelligent SOPs.

        Args:
            playbook_content: Playbook definition (JSON, YAML, etc.)
            playbook_format: Source SOAR platform format
            uploaded_by: Username of uploader

        Returns:
            Processing result with created entries
        """

        api_key = os.getenv('ANTHROPIC_API_KEY')
        if not api_key:
            logger.warning("ANTHROPIC_API_KEY not set for SOAR conversion")
            return {"error": "AI service not configured for playbook conversion"}

        prompt = f"""You are a SOAR playbook analyst. Convert this playbook from {playbook_format} format into structured SOC knowledge base entries.

PLAYBOOK:
{playbook_content[:20000]}

Analyze this playbook and create knowledge base entries that capture:
1. The workflow/procedure as a human-readable SOP
2. Decision points and conditions
3. Actions to take at each step
4. Escalation criteria
5. Expected outcomes

Convert to JSON:
{{
  "playbook_summary": "What this playbook does",
  "original_format": "{playbook_format}",
  "entries": [
    {{
      "title": "Entry title",
      "content_type": "playbook",
      "category": "incident_response",
      "content": "Full step-by-step procedure",
      "tags": ["soar-converted", "{playbook_format}"],
      "severity_filter": ["high", "critical"],
      "incident_types": ["type"],
      "priority": 50,
      "ai_summary": "Brief summary",
      "extracted_rules": ["Step 1", "Step 2"]
    }}
  ]
}}"""

        try:
            import aiohttp

            headers = {
                "x-api-key": api_key,
                "anthropic-version": "2023-06-01",
                "content-type": "application/json"
            }

            payload = {
                "model": self.model,
                "max_tokens": 4000,
                "messages": [{"role": "user", "content": prompt}]
            }

            async with aiohttp.ClientSession() as session:
                async with session.post(
                    "https://api.anthropic.com/v1/messages",
                    headers=headers,
                    json=payload,
                    timeout=aiohttp.ClientTimeout(total=60)
                ) as response:
                    if response.status != 200:
                        error_text = await response.text()
                        return {"error": f"AI API error: {error_text[:500]}"}

                    data = await response.json()
                    response_text = data['content'][0]['text']

                    analysis = self._parse_analysis_response(response_text)

                    if analysis.get('error'):
                        return analysis

                    # Create entries
                    created_entries = await self._create_entries_from_analysis(
                        analysis,
                        f"soar_import_{playbook_format}",
                        playbook_format,
                        uploaded_by
                    )

                    return {
                        "status": "completed",
                        "playbook_format": playbook_format,
                        "analysis": analysis,
                        "created_entries": created_entries,
                        "entries_count": len(created_entries)
                    }

        except Exception as e:
            logger.error(f"SOAR playbook processing failed: {e}")
            return {"error": str(e)}


# Singleton instance
kb_document_processor = KBDocumentProcessor()


def get_kb_document_processor() -> KBDocumentProcessor:
    """Get the singleton document processor instance."""
    return kb_document_processor
